import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

slides_data = [
    {
        "title": "Vigil-Core: Digital Forensics Platform",
        "content": [
            "Team: Digital Forensic",
            "Duration: 30 Minutes",
            "Team Size: 6 Members",
            "Focus: Malware and Network Forensic Analysis",
            "A unified platform for Security Operations Centers (SOC).",
            "Designed for high-fidelity detection and correlation of threats."
        ]
    },
    {
        "title": "Table of Contents",
        "content": [
            "1. Introduction to Digital Forensics",
            "2. Platform Architecture & Problems Solved",
            "3. Malware Team Overview (Arnab, Ankita, Karan)",
            "4. Malware Forensics Phases 1 to 5",
            "5. Network Team Overview (Cibin, Chetti Priyanka, Aditi)",
            "6. Network Forensics: Deep Packet Inspection & Narratives",
            "7. Demo Flow and System Performance",
            "8. Conclusion & Impact"
        ]
    },
    {
        "title": "Introduction to Digital Forensics",
        "content": [
            "Digital forensics is essential for incident response.",
            "Malware Forensics: Involves analyzing malicious binaries and memory.",
            "Network Forensics: Analyzes lateral movement and data exfiltration.",
            "Vigil-Core provides a 'Single Pane of Glass'.",
            "Reduces manual triage time drastically through automation.",
            "Integrates dynamic analysis into a cohesive dashboard."
        ]
    },
    {
        "title": "Platform Infrastructure & Stack",
        "content": [
            "Backend: Python 3, Flask, Volatility 3, YARA.",
            "Frontend: React, Vite, Tailwind CSS, Recharts.",
            "State & Streaming: SQLite for chunked memory streaming.",
            "Caching: Redis for O(1) Threat Intel caching.",
            "Network Sensors: Zeek and Suricata rulesets.",
            "Highly modular micro-architecture allows parallel development."
        ]
    },
    {
        "title": "Challenges: RAM Exhaustion",
        "content": [
            "Problem: Processing 20GB+ PCAP files crashed the backend.",
            "Network traffic telemetry is massive and cannot fit in memory.",
            "Solution: We engineered a chunked SQLite streaming engine.",
            "Data is parsed in small megabyte fragments and written to disk.",
            "This ensures constant, low memory overhead during analysis.",
            "Enables enterprise-scale forensic processing on consumer hardware."
        ]
    },
    {
        "title": "Challenges: Rate Limits & Freezing",
        "content": [
            "Problem: Querying VirusTotal for every IP led to API rate limiting.",
            "Solution: Built an O(1) Redis caching layer to remember known IPs.",
            "Problem: Analysis blocked the main UI thread.",
            "Solution: We utilized asynchronous Python subprocesses.",
            "React dashboard queries status continuously without freezing.",
            "Ensures flawless user experience during heavy workloads."
        ]
    },
    {
        "title": "Malware Team Overview",
        "content": [
            "Team Leader: Arnab",
            "Members: Ankita, Karan",
            "Focus: Building the CyArt Malware Engine.",
            "Goal: Create a pipeline to handle memory dumps sequentially.",
            "Scope covers everything from file validation to report generation.",
            "Divided the complexity into 5 distinct processing phases."
        ]
    },
    {
        "title": "Arnab: CyArt Malware Intro",
        "content": [
            "Speaker: Arnab (Team Leader)",
            "Engineered the foundational CyArt Malware backend.",
            "Integrated Volatility 3 for raw memory extraction.",
            "Designed the asynchronous execution loop for Flask.",
            "Ensures zero data loss while streaming terminal logs to the UI.",
            "Established the baseline framework for the 5-phase extraction."
        ]
    },
    {
        "title": "Arnab: Malware Architecture Flow",
        "content": [
            "Speaker: Arnab",
            "The upload initiates the pipeline asynchronously.",
            "Phase 1: Validates the file and calculates cryptographic hashes.",
            "Phase 2: Scans for encrypted/packed regions using entropy.",
            "Phase 3: Carves out executable payloads natively.",
            "Phase 4: Extracts Indicators of Compromise (IOCs).",
            "Phase 5: Generates final actionable reports."
        ]
    },
    {
        "title": "Arnab: Memory Dump Samples",
        "content": [
            "Speaker: Arnab",
            "Tested against real-world artifacts in the /samples folder.",
            "Target: memdump.mem (536 MB).",
            "This file simulates a raw memory capture of an infected host.",
            "We analyze the memory map to find injected payloads.",
            "Provides a high-fidelity test case for Volatility 3 extraction."
        ]
    },
    {
        "title": "Ankita: Phase 1 (Dump Validator)",
        "content": [
            "Speaker: Ankita",
            "I developed the initial ingestion phases for the malware pipeline.",
            "Phase 1 validates the integrity of the uploaded forensic file.",
            "Verifies file formats (e.g., .mem, .exe, .dll, .pcap).",
            "Computes SHA-256 and MD5 cryptographic hashes.",
            "Maintains chain of custody for digital evidence logging."
        ]
    },
    {
        "title": "Ankita: Phase 1 Deep Dive",
        "content": [
            "Speaker: Ankita",
            "If a file is corrupted, the pipeline fails safely immediately.",
            "The computed hash is instantly queried against our Redis cache.",
            "If it matches a known malicious signature, it alerts the SOC.",
            "Reduces redundant scanning of previously analyzed files.",
            "Sets the baseline parameters for the Volatility engine."
        ]
    },
    {
        "title": "Ankita: Phase 2 (Region Scanner)",
        "content": [
            "Speaker: Ankita",
            "Once validated, the file enters Phase 2 for region scanning.",
            "The engine analyzes data entropy block by block.",
            "We established an entropy threshold of 7.5 (out of 8.0).",
            "High entropy indicates packed, compressed, or encrypted code.",
            "Identifies hidden malicious logic within seemingly benign files."
        ]
    },
    {
        "title": "Ankita: Phase 2 Visualization",
        "content": [
            "Speaker: Ankita",
            "The output of Phase 2 is mapped and fed to the UI.",
            "Generates visual representations of memory allocations.",
            "Highlights the exact memory offsets containing threats.",
            "Allows incident responders to focus solely on encrypted regions.",
            "Significantly narrows down the scope for payload extraction."
        ]
    },
    {
        "title": "Karan: Phase 3 (Payload Extractor)",
        "content": [
            "Speaker: Karan",
            "I engineered the extraction and reporting modules.",
            "Phase 3 takes the flagged offsets from Ankita's region scanner.",
            "It actively carves raw executables natively from memory blocks.",
            "Extracts embedded shellcode and isolates it for testing.",
            "Validates Portable Executable (PE) headers for carved binaries."
        ]
    },
    {
        "title": "Karan: Phase 3 Deep Dive",
        "content": [
            "Speaker: Karan",
            "Without Phase 3, we would only know a threat exists.",
            "By carving the executable, we can reverse-engineer it.",
            "Outputs the extracted binaries securely to an isolated folder.",
            "Ensures no accidental execution occurs on the analyst's machine.",
            "Prepares the raw payload for deep static signature scanning."
        ]
    },
    {
        "title": "Karan: Phase 4 (IOC Extractor)",
        "content": [
            "Speaker: Karan",
            "Phase 4 scans the carved payloads against threat intelligence.",
            "Executes massive enterprise YARA rule sets against the files.",
            "Extracts embedded IPs, domains, and malicious URLs.",
            "Runs string pattern analysis to find hidden command instructions.",
            "Cross-references extracted domains with AlienVault OTX."
        ]
    },
    {
        "title": "Karan: Phase 4 & Redis Caching",
        "content": [
            "Speaker: Karan",
            "We integrated the O(1) Redis caching layer heavily here.",
            "Prevents our threat intelligence API keys from being exhausted.",
            "Ensures IOC extraction occurs in milliseconds instead of minutes.",
            "Flags known Advanced Persistent Threat (APT) indicators.",
            "All findings are structured into a standardized JSON format."
        ]
    },
    {
        "title": "Karan: Phase 5 (Report Generation)",
        "content": [
            "Speaker: Karan",
            "The final phase compiles all previous outputs natively.",
            "Generates an automated Executive Summary for management.",
            "Compiles JSON / PDF summaries for archival.",
            "Prepares Splunk-compatible payloads for SIEM delivery.",
            "Hydrates the React UI data components with real-time metrics."
        ]
    },
    {
        "title": "Network Team Overview",
        "content": [
            "Team Leader: Cibin",
            "Members: Chetti Priyanka, Aditi",
            "Focus: Building the NetForensicX Pipeline.",
            "Goal: Process massive PCAP network captures safely.",
            "Designed the architecture around streaming data ingestion.",
            "Translates raw packets into an actionable narrative."
        ]
    },
    {
        "title": "Cibin: NetForensicX Intro",
        "content": [
            "Speaker: Cibin (Team Leader)",
            "Engineered the overarching network ingestion framework.",
            "We had to process gigabytes of data on limited hardware.",
            "Traditional tools load everything into RAM and crash.",
            "NetForensicX avoids this by writing states to SQLite immediately.",
            "Provides an unbreakable, sequential packet parsing layer."
        ]
    },
    {
        "title": "Cibin: Network Architecture Pipeline",
        "content": [
            "Speaker: Cibin",
            "Step 1: Upload massive PCAP to the Flask backend.",
            "Step 2: Stream packets into the Zeek/Suricata IDS engine.",
            "Step 3: Chunk processed logs directly into SQLite.",
            "Step 4: Execute Node Correlation across the database.",
            "Step 5: Generate the sequential Attack Story timeline."
        ]
    },
    {
        "title": "Chetti Priyanka: Deep Packet Inspection",
        "content": [
            "Speaker: Chetti Priyanka",
            "I integrated Zeek and Suricata for Deep Packet Inspection.",
            "DPI goes beyond IP/Port matching; it inspects packet payloads.",
            "Parses contents directly from the SQLite stream natively.",
            "Extracts JA3 SSL/TLS hashes to track encrypted connections.",
            "Identifies HTTP host header anomalies and spoofing."
        ]
    },
    {
        "title": "Chetti Priyanka: PCAP Samples",
        "content": [
            "Speaker: Chetti Priyanka",
            "We validated the pipeline using 3 real-world samples:",
            "Hive_06082021.pcap (57.1 MB): Analyzes Ransomware C2 Traffic.",
            "sf19us-MTA-lab-16.pcap (15.8 MB): Malware Traffic Analysis.",
            "amp.TCP.reflection.SYNACK.pcap: Analyzes DDoS reflections.",
            "Ensured robust detection across wildly different attack vectors."
        ]
    },
    {
        "title": "Chetti Priyanka: Threat Signatures",
        "content": [
            "Speaker: Chetti Priyanka",
            "Our Suricata rulesets flag specific malicious behaviors.",
            "Detects multi-protocol brute-force attacks natively.",
            "Identifies cleartext protocol usage (Telnet, FTP).",
            "Generates alerts that are instantly routed to the UI timeline.",
            "Crucial for identifying the initial access vector in a breach."
        ]
    },
    {
        "title": "Aditi: Attack Narrative Generator",
        "content": [
            "Speaker: Aditi",
            "I developed the Attack Narrative Generation logic.",
            "Instead of showing analysts 10,000 raw, unreadable packets...",
            "We correlate events into a Cyber Kill Chain.",
            "Reconnaissance -> Initial Access -> C2 -> Exfiltration.",
            "Translates raw network logs into a human-readable story."
        ]
    },
    {
        "title": "Aditi: Correlation Logic",
        "content": [
            "Speaker: Aditi",
            "The correlation logic groups disparate events by timestamp.",
            "If IP A scans IP B, and then IP A transfers a payload...",
            "The system binds these two events into a single threat timeline.",
            "Saves SOC teams hours of manual packet tracing in Wireshark.",
            "Allows junior analysts to understand complex network breaches."
        ]
    },
    {
        "title": "Aditi: UI Data Visualization",
        "content": [
            "Speaker: Aditi",
            "We visualize this network telemetry directly in the dashboard.",
            "Integrated React Recharts for high-performance rendering.",
            "Protocol Distribution Pie Chart: Splits TCP/UDP/ICMP.",
            "Service Distribution Bar Graph: Shows heavily targeted hosts.",
            "Nodes are mapped visually to highlight compromised endpoints."
        ]
    },
    {
        "title": "Demo Flow: Malware Analysis",
        "content": [
            "Step 1: Authenticate and hit the dashboard.",
            "Step 2: Upload memdump.mem.",
            "Step 3: Observe asynchronous Python sub-processes firing.",
            "UI updates instantly as Phases 1 through 5 complete.",
            "Results display the Threat Level Banner in High-Alert Red.",
            "YARA matches and IOCs are dynamically tabulated."
        ]
    },
    {
        "title": "Demo Flow: Network Analysis",
        "content": [
            "Step 1: Navigate to the Network Forensics module.",
            "Step 2: Upload Hive_06082021.pcap.",
            "The Attack Story populates the right-hand panel instantly.",
            "Protocol / Service graphs render dynamically.",
            "Compromised hosts are tabulated in real-time.",
            "Demonstrates the power of the SQLite streaming engine."
        ]
    },
    {
        "title": "Demo Flow: History & Persistence",
        "content": [
            "Navigate to 'Report History / Saved Scans'.",
            "Click on a historical analysis report from a prior session.",
            "Data persists flawlessly via the SQLite state backend.",
            "No re-scanning is ever required for past forensic events.",
            "Users can modify their profiles and settings without losing data."
        ]
    },
    {
        "title": "Performance Metrics",
        "content": [
            "Dashboard Latency: ~4.3ms load times.",
            "Processing Speed: 20GB PCAP parsed in under 3 minutes.",
            "API Optimization: Near 100% Redis cache hit rate for common IPs.",
            "Reliability: Zero memory crashes during heavy testing.",
            "Scalability: Designed for enterprise SOC environments."
        ]
    },
    {
        "title": "Platform Impact & Benefits",
        "content": [
            "Replaces 5 disjointed terminal tools with 1 unified platform.",
            "Junior Analysts can read plain-text 'Attack Stories' easily.",
            "Reduces Incident Response Mean Time To Respond (MTTR).",
            "Provides a scalable template for adding new forensic tools.",
            "Automates the most tedious aspects of digital forensics."
        ]
    },
    {
        "title": "Conclusion",
        "content": [
            "Vigil-Core bridges technical operations and SOC management.",
            "Divided complexity into modular Python engines.",
            "Successfully integrated React for high-performance visualizations.",
            "Overcame massive data ingestion challenges natively.",
            "The platform is stable, functional, and ready for deployment."
        ]
    },
    {
        "title": "Questions?",
        "content": [
            "Thank you for your time.",
            "We will now open the floor for questions.",
            "Digital Forensic Team:",
            "- Malware: Arnab, Ankita, Karan",
            "- Network: Cibin, Chetti Priyanka, Aditi",
            "Project: Vigil-Core"
        ]
    }
]

for sd in slides_data:
    slide_layout = prs.slide_layouts[1] # Title and Content Layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Configure Title
    title = slide.shapes.title
    title.text = sd["title"]
    for paragraph in title.text_frame.paragraphs:
        # Cannot easily enforce Montserrat without system font, but will set name
        paragraph.font.name = 'Montserrat'
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0x38, 0xBD, 0xF8)
    
    # Configure Body
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    if sd["content"]:
        tf.text = sd["content"][0]
        tf.paragraphs[0].font.name = 'Inter'
        tf.paragraphs[0].font.size = Pt(20)
        
        for point in sd["content"][1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.name = 'Inter'
            p.font.size = Pt(20)

prs.save("VigilCore_Presentation.pptx")
print("Saved VigilCore_Presentation.pptx successfully!")
